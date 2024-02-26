import os
import shutil 
from docx import Document
from pptx import Presentation
from googletrans import Translator

def translate_vtt(input_file, output_file):
    with open(input_file, 'r', encoding='utf-8') as file:
        content = file.read()

    sentences = content.split('. ')

    translator = Translator()

    translated_sentences = []
    for sentence in sentences:
        translated_sentence = translator.translate(sentence, src='en', dest='ro').text
        translated_sentences.append(translated_sentence)

    translated_content = '. '.join(translated_sentences)

    with open(output_file, 'w', encoding='utf-8') as file:
        file.write(translated_content)

def translate_docx(input_file, output_file):
    if not os.path.exists(input_file):
        print("Fișierul de intrare nu există.")
        return

    doc_copy = input_file.replace(".docx", "_Ro.docx")
    shutil.copy(input_file, doc_copy)

    doc = Document(doc_copy)

    translator = Translator()

    for paragraph in doc.paragraphs:
        translated_text = translator.translate(paragraph.text, src='en', dest='ro').text
        paragraph.text = translated_text

    doc.save(output_file)
    print(f"Documentul a fost tradus și salvat ca {output_file}.")

def translate_pptx(input_file, output_file):
    if not os.path.exists(input_file):
        print("Fișierul de intrare nu există.")
        return

    pptx_copy = input_file.replace(".pptx", "_translated.pptx")
    shutil.copy(input_file, pptx_copy)

    prs = Presentation(pptx_copy)
    translator = Translator()

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                try:
                    translated_text = translator.translate(shape.text, src='en', dest='ro').text
                    shape.text = translated_text
                except Exception as e:
                    print(f"Eroare la traducerea textului: {e}")
                    print("Textul nu a putut fi tradus și va rămâne neschimbat.")

    prs.save(output_file)
    print(f"Prezentarea PowerPoint a fost tradusă și salvată ca {output_file}.")



def translate_documents_in_folder(folder_path):
    for file in os.listdir(folder_path):
        if file.endswith(".docx"):
            input_file = os.path.join(folder_path, file)
            output_file = os.path.join(folder_path, f"_translated{file}")
            translate_docx(input_file, output_file)
        elif file.endswith(".vtt"):
            input_file = os.path.join(folder_path, file)
            output_file = os.path.join(folder_path, f"{file}_Ro")
            translate_vtt(input_file, output_file)
        elif file.endswith(".pptx"):
            input_file = os.path.join(folder_path, file)
            output_file = os.path.join(folder_path, f"_translated{file}")
            translate_pptx(input_file, output_file)
        else:
            print(f"Ignorând {file} - format necunoscut.")


translate_documents_in_folder(r"C:\Users\danie\Desktop\munca\C l Green Entrepreneurship -20230803T052345Z-001\C l Green Entrepreneurship\0.Intro")
